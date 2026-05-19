#!/usr/bin/env python3
"""Build a hackathon pitch deck for Repo Health Intelligence."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

BG = RGBColor(0x0A, 0x0A, 0x0C)
PANEL = RGBColor(0x18, 0x18, 0x1B)
ACCENT = RGBColor(0x22, 0xC5, 0x5E)
ACCENT_DIM = RGBColor(0x16, 0x99, 0x4A)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
TEXT = RGBColor(0xE8, 0xE8, 0xEA)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
BORDER = RGBColor(0x27, 0x27, 0x2A)

W, H = Inches(13.333), Inches(7.5)

def new_deck():
    p = Presentation()
    p.slide_width = W
    p.slide_height = H
    return p

def slide_blank(p):
    layout = p.slide_layouts[6]  # blank
    s = p.slides.add_slide(layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font="Helvetica Neue"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        runs = [text]
    else:
        runs = text
    for i, t in enumerate(runs):
        run = p.add_run() if i > 0 else p.runs[0] if p.runs else p.add_run()
        run.text = t
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + it
        run.font.name = "Helvetica Neue"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return tb

def panel(slide, x, y, w, h, fill=PANEL, border=BORDER):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(0.75)
    sh.adjustments[0] = 0.05
    sh.shadow.inherit = False
    return sh

def header(slide, title, kicker=None):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
             "REPO HEALTH INTELLIGENCE" if kicker is None else kicker,
             size=10, color=ACCENT, bold=True)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=TEXT)

def footer(slide, idx, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             f"github.com/avinrique/repo-health-intelligence", size=10, color=MUTED)
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.3),
             f"{idx} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def kpi(slide, x, y, w, h, value, label, color=ACCENT):
    panel(slide, x, y, w, h)
    add_text(slide, x, y + Inches(0.15), w, Inches(0.7), value, size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.95), w, Inches(0.3), label, size=10, color=MUTED, align=PP_ALIGN.CENTER)

def build():
    p = new_deck()
    slides_meta = []

    # 1 TITLE
    s = slide_blank(p)
    add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4), "HACKATHON SUBMISSION · 2026", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.2), "Repo Health Intelligence", size=58, bold=True, color=TEXT)
    add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(0.6),
             "Track how a codebase evolves over time — and tell teams whether they're winning or losing.",
             size=20, color=MUTED)
    panel(s, Inches(0.6), Inches(3.9), Inches(12.1), Inches(2.2))
    add_text(s, Inches(0.95), Inches(4.1), Inches(11.4), Inches(0.4), "WHAT IT DOES", size=10, color=ACCENT, bold=True)
    add_bullets(s, Inches(0.95), Inches(4.45), Inches(11.4), Inches(1.6), [
        "Ingests a public git repo and builds a per-commit code knowledge graph with tree-sitter",
        "Composite health score with 5 explainable subscores — drift, tests, hotspots, deps, arch",
        "Click any dip on the time-series → LLM tells you why",
        "Pre-merge prediction, bus factor, architectural drift, force-directed graph viz",
    ])
    add_text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4), "Tested on tarinagarwal/Edulume — 109 commits ingested in ~30s on an M2", size=12, color=MUTED)
    footer(s, 1, 13)

    # 2 PROBLEM
    s = slide_blank(p)
    header(s, "The problem")
    add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.8),
             "Linters and coverage reports describe a codebase today.",
             size=24, color=TEXT)
    add_text(s, Inches(0.6), Inches(2.3), Inches(12), Inches(0.8),
             "No one connects \"this PR merged\" to \"and here's how it shifted the complexity,",
             size=24, color=TEXT)
    add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(0.8),
             "architecture, and risk of the codebase.\"",
             size=24, color=TEXT)
    add_text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.8), "Engineering teams ship blind.", size=28, color=AMBER, bold=True)
    panel(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.6))
    add_text(s, Inches(0.95), Inches(5.15), Inches(11.4), Inches(0.4), "THE GAP", size=10, color=ACCENT, bold=True)
    add_bullets(s, Inches(0.95), Inches(5.5), Inches(11.4), Inches(1.0), [
        "Coverage reports = static snapshot of one moment",
        "Git history is rich, but no tooling converts it into a longitudinal health signal",
        "Teams notice rot only when something breaks",
    ])
    footer(s, 2, 13)

    # 3 SOLUTION OVERVIEW
    s = slide_blank(p)
    header(s, "How it works")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.5), "Four-stage pipeline · zero LLM dependency for the core metrics", size=14, color=MUTED)
    # 4 columns
    titles = ["1 · INGEST", "2 · GRAPH", "3 · SCORE", "4 · DASHBOARD"]
    bodies = [
        ["Clone repo", "Walk every commit", "tree-sitter parses TS/TSX/JS/JSX/Py", "Sample for repos > 500"],
        ["Per-commit nodes (files)", "Edges = resolved imports", "Tarjan SCC for cycles", "Stored in SQLite"],
        ["5 normalized subscores", "Composite 0–100", "Each call is explainable", "Baseline = first commit"],
        ["Time-series chart", "Click dip → LLM explains", "Bus factor heatmap", "Pre-merge simulator"],
    ]
    cx = Inches(0.6)
    cw = Inches(2.95)
    gap = Inches(0.15)
    for i in range(4):
        x = cx + (cw + gap) * i
        panel(s, x, Inches(2.5), cw, Inches(3.8))
        add_text(s, x + Inches(0.2), Inches(2.65), cw - Inches(0.4), Inches(0.4), titles[i], size=10, color=ACCENT, bold=True)
        add_bullets(s, x + Inches(0.2), Inches(3.05), cw - Inches(0.4), Inches(3.2), bodies[i], size=12)
    footer(s, 3, 13)

    # 4 HEALTH SCORE FORMULA
    s = slide_blank(p)
    header(s, "The health score")
    add_text(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.5), "Five explainable subscores → composite 0–100", size=14, color=MUTED)
    panel(s, Inches(0.6), Inches(2.35), Inches(12.1), Inches(1.05), fill=RGBColor(0x10, 0x10, 0x12))
    add_text(s, Inches(0.85), Inches(2.55), Inches(11.7), Inches(0.65),
             "health = 25·(1 − drift) + 25·tests + 20·(1 − hotspots) + 15·(1 − rot) + 15·(1 − arch_drift)",
             size=18, color=ACCENT, font="Menlo")
    # subscore cards
    subs = [
        ("Complexity drift", "avg cyclomatic complexity vs first commit", AMBER),
        ("Test coverage", "test files / source files (path heuristic — no run needed)", ACCENT),
        ("Hotspot risk", "top-10 concentration + log-intensity of churn × cx", RED),
        ("Dependency rot", "num deps + commits since package.json last bumped", RGBColor(0xA7, 0x8B, 0xFA)),
        ("Arch drift", "growth in cycles + orphan ratio + max fan-in", RGBColor(0xF4, 0x72, 0xB6)),
    ]
    cy = Inches(3.7)
    ch = Inches(2.95)
    cw = (W - Inches(0.6) - Inches(0.6) - Inches(0.5)) / 5
    cx = Inches(0.6)
    gap = Inches(0.125)
    for i, (n, body, c) in enumerate(subs):
        x = cx + (cw + gap) * i
        panel(s, x, cy, cw, ch)
        # accent bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.18), cy + Inches(0.18), Inches(0.06), ch - Inches(0.36))
        bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = c
        add_text(s, x + Inches(0.35), cy + Inches(0.2), cw - Inches(0.5), Inches(0.5), n, size=14, bold=True, color=TEXT)
        add_text(s, x + Inches(0.35), cy + Inches(0.65), cw - Inches(0.5), ch - Inches(0.8), body, size=11, color=MUTED)
    footer(s, 4, 13)

    # 5 EDULUME RESULTS
    s = slide_blank(p)
    header(s, "Live result: tarinagarwal/Edulume", kicker="DEMO RUN")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), "109 commits ingested · ~30s wall-clock on M2", size=13, color=MUTED)
    # KPIs row
    kpi_y = Inches(2.4)
    kpi_w = Inches(2.4)
    kpi_h = Inches(1.3)
    gap = Inches(0.15)
    items = [
        ("55.2", "Starting health", ACCENT),
        ("28.2", "Current health", RED),
        ("−27.0", "Net drop", AMBER),
        ("3", "Flagged dips", AMBER),
        ("30", "Single-owner files", RED),
    ]
    cx = Inches(0.6)
    for i, (v, l, c) in enumerate(items):
        x = cx + (kpi_w + gap) * i
        kpi(s, x, kpi_y, kpi_w, kpi_h, v, l, color=c)
    # findings
    panel(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.7))
    add_text(s, Inches(0.95), Inches(4.25), Inches(11.4), Inches(0.4), "WHAT THE SYSTEM CAUGHT", size=10, color=ACCENT, bold=True)
    add_bullets(s, Inches(0.95), Inches(4.65), Inches(11.4), Inches(2.0), [
        "Biggest dip at commit #1 \"add discussion forum\" — health drops 18 pts; LLM explains: drift + new hotspots",
        "client/src/utils/api.ts is the #1 hotspot — churn 19 × complexity 54 = risk 1026",
        "Real circular import: TestTimeRemaining.tsx ↔ TestPageStandalone.tsx (caught by Tarjan SCC)",
        "tarinagarwal owns 87% of server/routes/ and 100% of auth code — bus factor = 1",
        "Simulating feature/olive-theme branch predicts a −2.1 health drop, mostly from hotspot risk",
    ], size=13)
    footer(s, 5, 13)

    # 6 TIME-SERIES + DIPS
    s = slide_blank(p)
    header(s, "Time-series dashboard")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), "Health trend, dip annotations, click-to-explain", size=14, color=MUTED)
    # mock chart panel
    panel(s, Inches(0.6), Inches(2.4), Inches(7.6), Inches(4.3))
    add_text(s, Inches(0.85), Inches(2.55), Inches(7), Inches(0.4), "Health (0–100)", size=10, color=MUTED, bold=True)
    # draw fake polyline using a sequence of small rectangles - skip; just label
    # axis
    axis = s.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(1.0), Inches(6.2), Inches(7.0), Emu(0)) if False else None
    # draw line approximation using a freeform: use small lines via add_connector
    pts = [(0,55),(8,42),(15,40),(22,32),(28,30),(36,28),(45,27),(60,26),(75,29),(95,28),(108,28)]
    x0 = Inches(1.0); y0 = Inches(6.0)
    xspan = Inches(7.0); yspan = Inches(3.0)
    # baseline
    base = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y0, xspan, Emu(8000))
    base.fill.solid(); base.fill.fore_color.rgb = BORDER; base.line.fill.background()
    # plot polyline as connectors
    from pptx.enum.shapes import MSO_CONNECTOR
    def pt_to_emu(i, v):
        x = x0 + Emu(int(xspan / 108 * i))
        y = y0 - Emu(int(yspan / 60 * v))
        return x, y
    for a, b in zip(pts, pts[1:]):
        ax, ay = pt_to_emu(*a); bx, by = pt_to_emu(*b)
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax, ay, bx, by)
        ln.line.color.rgb = ACCENT; ln.line.width = Pt(2.5)
    # dip dots
    for dx, dy in [(8,42),(28,30),(60,26)]:
        x, y = pt_to_emu(dx, dy)
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.08), y - Inches(0.08), Inches(0.16), Inches(0.16))
        d.fill.solid(); d.fill.fore_color.rgb = RED
        d.line.color.rgb = TEXT; d.line.width = Pt(0.5)
    # right column callout
    panel(s, Inches(8.4), Inches(2.4), Inches(4.3), Inches(4.3))
    add_text(s, Inches(8.65), Inches(2.55), Inches(4), Inches(0.4), "CLICK A RED DOT", size=10, color=ACCENT, bold=True)
    add_text(s, Inches(8.65), Inches(2.9), Inches(4), Inches(0.5), "Commit #1: add discussion forum", size=15, bold=True, color=TEXT)
    add_text(s, Inches(8.65), Inches(3.3), Inches(4), Inches(0.3), "10a91b6 — Aug 20, 2025", size=10, color=MUTED, font="Menlo")
    panel(s, Inches(8.65), Inches(3.75), Inches(3.85), Inches(2.8), fill=RGBColor(0x10, 0x10, 0x12))
    add_text(s, Inches(8.85), Inches(3.9), Inches(3.5), Inches(0.35), "WHY THE DROP", size=9, color=AMBER, bold=True)
    add_text(s, Inches(8.85), Inches(4.2), Inches(3.5), Inches(2.4),
             "Health fell 18.0 points; the dominant driver was complexity drift worsened. "
             "Largest swings: server/db.js (Δcx=+7), Navbar.tsx (Δcx=+2). New hotspots: DiscussionDetailPage.tsx, routes/discussions.js. "
             "Files 27→35, LOC 3587→7097.",
             size=11, color=TEXT)
    footer(s, 6, 13)

    # 7 KNOWLEDGE GRAPH DIFF
    s = slide_blank(p)
    header(s, "Knowledge graph diff")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), "Compare any two commits. See exactly what changed structurally.", size=14, color=MUTED)
    # boxes for added / removed / changed
    cols = [
        ("+ 8", "files added", ACCENT, ["DiscussionDetailPage.tsx", "MentionInput.tsx", "NotificationDropdown.tsx", "CreateDiscussionPage.tsx", "DiscussionsPage.tsx"]),
        ("− 0", "files removed", RGBColor(0x6B, 0x72, 0x80), []),
        ("~ 11", "files changed", AMBER, ["server/db.js  Δcx +7", "Navbar.tsx  Δcx +2", "client/utils/api.ts  Δcx +1"]),
        ("+ 22", "edges added", ACCENT, ["DiscussionsPage → api.ts", "DiscussionDetail → MentionInput", "Navbar → NotificationDropdown"]),
    ]
    cw = Inches(2.95); ch = Inches(4.2); cx = Inches(0.6); cy = Inches(2.4); gap = Inches(0.15)
    for i, (val, lbl, c, items) in enumerate(cols):
        x = cx + (cw + gap) * i
        panel(s, x, cy, cw, ch)
        add_text(s, x, cy + Inches(0.2), cw, Inches(0.7), val, size=36, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.95), cw, Inches(0.3), lbl, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        if items:
            add_bullets(s, x + Inches(0.25), cy + Inches(1.4), cw - Inches(0.5), ch - Inches(1.5), items, size=10, color=TEXT)
    footer(s, 7, 13)

    # 8 HOTSPOT MAP
    s = slide_blank(p)
    header(s, "Hotspot map — where bugs live")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), "churn × complexity, ranked. Top 5 at Edulume HEAD:", size=14, color=MUTED)
    hot = [
        ("client/src/utils/api.ts", 1026, 19, 54),
        ("server/routes/auth.js", 836, 11, 76),
        ("server/routes/discussions.js", 786, 6, 131),
        ("client/src/components/layout/Navbar.tsx", 620, 10, 62),
        ("client/src/components/ai/PdfChatbotPage.tsx", 408, 6, 68),
    ]
    maxr = hot[0][1]
    y = Inches(2.5)
    rh = Inches(0.8)
    for path, risk, churn, cx in hot:
        panel(s, Inches(0.6), y, Inches(12.1), rh - Inches(0.1))
        add_text(s, Inches(0.85), y + Inches(0.1), Inches(7), Inches(0.4), path, size=13, bold=True, color=TEXT, font="Menlo")
        add_text(s, Inches(0.85), y + Inches(0.4), Inches(7), Inches(0.3),
                 f"churn {churn} · complexity {cx} · risk {risk}", size=11, color=MUTED)
        # bar
        pct = risk / maxr
        full_w = Inches(3.6)
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), y + Inches(0.3), full_w, Inches(0.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = BORDER; bar.line.fill.background()
        fg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), y + Inches(0.3), Inches(3.6 * pct), Inches(0.15))
        c = RED if pct > 0.66 else AMBER if pct > 0.33 else ACCENT
        fg.fill.solid(); fg.fill.fore_color.rgb = c; fg.line.fill.background()
        y += rh
    footer(s, 8, 13)

    # 9 BUS FACTOR
    s = slide_blank(p)
    header(s, "Bus factor — who owns what")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), "The smallest number of authors who together own > 50 % of commits on a file or module.", size=14, color=MUTED)
    mods = [
        ("client/src/", 25, "tarinagarwal", 73, 39),
        ("server/routes/", 6, "tarinagarwal", 87, 7),
        ("python-backend/fileUpload/", 2, "tarinagarwal", 100, 2),
        ("server/middleware/", 1, "tarinagarwal", 100, 1),
    ]
    y = Inches(2.4)
    rh = Inches(0.85)
    for mod, lowbus, top, share, files in mods:
        panel(s, Inches(0.6), y, Inches(12.1), rh - Inches(0.1))
        add_text(s, Inches(0.85), y + Inches(0.1), Inches(5), Inches(0.4), mod, size=14, bold=True, color=TEXT, font="Menlo")
        add_text(s, Inches(0.85), y + Inches(0.42), Inches(5), Inches(0.3),
                 f"{files} files · {lowbus} single-owner files", size=11, color=MUTED)
        # owner share bar
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), y + Inches(0.28), Inches(4.5), Inches(0.2))
        bar.fill.solid(); bar.fill.fore_color.rgb = BORDER; bar.line.fill.background()
        fg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), y + Inches(0.28), Inches(4.5 * share / 100), Inches(0.2))
        fg.fill.solid(); fg.fill.fore_color.rgb = RED if share >= 80 else AMBER; fg.line.fill.background()
        add_text(s, Inches(11.2), y + Inches(0.15), Inches(1.4), Inches(0.4), f"{top} {share}%", size=12, color=TEXT, align=PP_ALIGN.RIGHT)
        y += rh
    add_text(s, Inches(0.6), Inches(6.05), Inches(12), Inches(0.4), "Verdict: lose one person → most of Edulume needs new ownership.", size=14, color=AMBER, bold=True)
    footer(s, 9, 13)

    # 10 ARCH DRIFT
    s = slide_blank(p)
    header(s, "Architectural drift")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), "Cycles + orphans + max fan-in, baselined to the first commit.", size=14, color=MUTED)
    # grid: orphans/cycles/fanin at first vs HEAD
    rows = [
        ("Cycles (Tarjan SCC ≥2)", "0", "1", "TestTimeRemaining.tsx ↔ TestPageStandalone.tsx"),
        ("Orphan files", "6", "18", "files with no in/out internal imports — unused or dead"),
        ("Max fan-in", "~4", "~5", "most-imported file = api.ts utility"),
        ("Single-owner files", "0", "30", "every one is a bus-factor-1 risk"),
    ]
    y = Inches(2.5)
    rh = Inches(0.85)
    # header row
    panel(s, Inches(0.6), y, Inches(12.1), Inches(0.5), fill=BORDER)
    add_text(s, Inches(0.85), y + Inches(0.13), Inches(5), Inches(0.3), "METRIC", size=10, color=MUTED, bold=True)
    add_text(s, Inches(6.0), y + Inches(0.13), Inches(1.5), Inches(0.3), "BASELINE", size=10, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.8), y + Inches(0.13), Inches(1.5), Inches(0.3), "AT HEAD", size=10, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.5), y + Inches(0.13), Inches(3.2), Inches(0.3), "NOTE", size=10, color=MUTED, bold=True)
    y += Inches(0.55)
    for metric, b, head, note in rows:
        panel(s, Inches(0.6), y, Inches(12.1), rh - Inches(0.1))
        add_text(s, Inches(0.85), y + Inches(0.18), Inches(5), Inches(0.4), metric, size=12, color=TEXT, bold=True)
        add_text(s, Inches(6.0), y + Inches(0.18), Inches(1.5), Inches(0.4), b, size=18, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(7.8), y + Inches(0.18), Inches(1.5), Inches(0.4), head, size=18, color=RED if head != b else MUTED, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(9.5), y + Inches(0.22), Inches(3.2), Inches(0.5), note, size=10, color=MUTED)
        y += rh
    footer(s, 10, 13)

    # 11 PRE-MERGE PREDICTION
    s = slide_blank(p)
    header(s, "Pre-merge health prediction")
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4), "Pick any branch on the upstream repo, simulate it against main, get predicted Δhealth.", size=14, color=MUTED)
    # big numbers
    kpi(s, Inches(0.6), Inches(2.4), Inches(3.9), Inches(1.5), "28.2", "Current health (main)", color=TEXT)
    kpi(s, Inches(4.7), Inches(2.4), Inches(3.9), Inches(1.5), "26.1", "Predicted (feature/olive-theme)", color=AMBER)
    kpi(s, Inches(8.8), Inches(2.4), Inches(3.9), Inches(1.5), "− 2.1", "Δ if merged", color=RED)
    # subscore impact
    panel(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.7))
    add_text(s, Inches(0.95), Inches(4.25), Inches(11.4), Inches(0.4), "SUBSCORE IMPACT", size=10, color=ACCENT, bold=True)
    sub_lines = [
        ("complexity drift", "0.723", "0.752", True),
        ("test coverage", "0.000", "0.000", False),
        ("hotspot risk", "0.777", "0.845", True),
        ("dependency rot", "0.500", "0.500", False),
        ("arch drift", "0.380", "0.380", False),
    ]
    sy = Inches(4.7)
    for name, before, after, worse in sub_lines:
        add_text(s, Inches(0.95), sy, Inches(3.5), Inches(0.35), name, size=12, color=TEXT)
        add_text(s, Inches(4.6), sy, Inches(1.5), Inches(0.35), before, size=12, color=MUTED, font="Menlo")
        add_text(s, Inches(6.4), sy, Inches(0.5), Inches(0.35), "→", size=12, color=MUTED)
        add_text(s, Inches(7.0), sy, Inches(1.5), Inches(0.35), after, size=12, color=RED if worse else MUTED, bold=worse, font="Menlo")
        add_text(s, Inches(9.0), sy, Inches(3.5), Inches(0.35),
                 "worse — hotspot risk grows" if name == "hotspot risk" else
                 "worse — avg cx rises" if name == "complexity drift" else
                 "no change", size=11, color=AMBER if worse else MUTED)
        sy += Inches(0.4)
    footer(s, 11, 13)

    # 12 LLM COST DISCIPLINE
    s = slide_blank(p)
    header(s, "LLM usage — every call defended")
    panel(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.4))
    add_text(s, Inches(0.95), Inches(2.15), Inches(11.4), Inches(0.4), "POLICY", size=10, color=ACCENT, bold=True)
    add_bullets(s, Inches(0.95), Inches(2.55), Inches(11.4), Inches(4.0), [
        "The only LLM endpoint is /api/explain — narrates a dip when the user clicks one",
        "Result is cached in narratives(sha, prev_sha) — second click hits the cache",
        "Prompt body is ~600 tokens of pre-computed deltas — no source code, no per-function data",
        "Output capped at 280 tokens — Haiku 4.5 via Vercel AI Gateway",
        "When AI_GATEWAY_API_KEY is unset, a deterministic rule-based summary runs instead — feature still works at zero cost",
        "Cost scales with user clicks on dips, not repo size · ingestion does zero LLM work",
    ], size=14)
    footer(s, 12, 13)

    # 13 STACK + WRAP
    s = slide_blank(p)
    header(s, "Stack · checklist · ask")
    # left: stack
    panel(s, Inches(0.6), Inches(1.9), Inches(6.0), Inches(4.5))
    add_text(s, Inches(0.85), Inches(2.05), Inches(5.5), Inches(0.4), "STACK", size=10, color=ACCENT, bold=True)
    add_bullets(s, Inches(0.85), Inches(2.45), Inches(5.5), Inches(3.8), [
        "Next.js 15 · React 19 · TypeScript",
        "tree-sitter (TS / TSX / JS / JSX / Python)",
        "simple-git for history walk",
        "better-sqlite3 for graph storage",
        "Recharts + custom canvas force-directed viz",
        "Vercel AI Gateway · claude-haiku-4-5",
    ], size=13)
    # right: checklist
    panel(s, Inches(6.75), Inches(1.9), Inches(5.95), Inches(4.5))
    add_text(s, Inches(7.0), Inches(2.05), Inches(5.5), Inches(0.4), "BRIEF CHECKLIST", size=10, color=ACCENT, bold=True)
    items_done = [
        "Repo ingestion + per-commit graphs",
        "5 explainable subscores (brief asked ≥ 4)",
        "Time-series dashboard + dip annotations",
        "Knowledge graph diff",
        "Hotspot map",
        "Bus factor per module  [bonus]",
        "Architectural drift  [bonus]",
        "LLM narrative  [bonus]",
        "Pre-merge prediction  [bonus]",
        "Multi-language  [bonus]",
    ]
    cy = Inches(2.45)
    for it in items_done:
        add_text(s, Inches(7.0), cy, Inches(0.4), Inches(0.3), "✓", size=14, color=ACCENT, bold=True)
        add_text(s, Inches(7.35), cy, Inches(5.2), Inches(0.3), it, size=12, color=TEXT)
        cy += Inches(0.32)
    # bottom: links
    add_text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
             "github.com/avinrique/repo-health-intelligence · demo: tarinagarwal/Edulume · built for the Repo Health Intelligence brief",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 13, 13)

    p.save("repo-health-intelligence.pptx")
    print("wrote repo-health-intelligence.pptx with 13 slides")

if __name__ == "__main__":
    build()
